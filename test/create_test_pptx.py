#!/usr/bin/env python3
"""
テスト用PPTXファイル作成スクリプト
日本語フォントのテストも含む
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_test_presentation():
    """テスト用のPPTXファイルを生成"""
    prs = Presentation()
    
    # スライド1: タイトルスライド
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Test Presentation for Translation"
    subtitle.text = "This is a test PowerPoint file with various text elements"
    
    # スライド2: 箇条書き
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Key Features"
    tf = content.text_frame
    tf.text = "First bullet point with important information"
    
    p = tf.add_paragraph()
    p.text = "Second bullet point with additional details"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Sub-point under second bullet"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Third main bullet point"
    p.level = 0
    
    # スライド3: テーブル
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトルを追加
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Data Table Example"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    
    # テーブルを追加
    rows = 3
    cols = 3
    left = Inches(2)
    top = Inches(2)
    width = Inches(6)
    height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # ヘッダー行
    table.cell(0, 0).text = "Category"
    table.cell(0, 1).text = "Value"
    table.cell(0, 2).text = "Status"
    
    # データ行
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "$1,000,000"
    table.cell(1, 2).text = "Achieved"
    
    table.cell(2, 0).text = "Profit"
    table.cell(2, 1).text = "$250,000"
    table.cell(2, 2).text = "In Progress"
    
    # スライド4: 複数のテキストボックス
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # タイトル
    txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "Multiple Text Elements"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # テキストボックス1
    txBox1 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3.5), Inches(2))
    tf1 = txBox1.text_frame
    tf1.text = "Left side content"
    p = tf1.add_paragraph()
    p.text = "This is important information that needs to be translated accurately."
    p.font.size = Pt(14)
    
    # テキストボックス2
    txBox2 = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(3.5), Inches(2))
    tf2 = txBox2.text_frame
    tf2.text = "Right side content"
    p = tf2.add_paragraph()
    p.text = "Additional details and context for better understanding."
    p.font.size = Pt(14)
    p.font.italic = True
    
    # スライド5: 様々なフォントスタイル
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Font Styles Test"
    tf = content.text_frame
    tf.clear()  # 既存のテキストをクリア
    
    # 通常のテキスト
    p = tf.add_paragraph()
    p.text = "Normal text style"
    p.font.size = Pt(18)
    
    # 太字
    p = tf.add_paragraph()
    p.text = "Bold text style"
    p.font.bold = True
    p.font.size = Pt(18)
    
    # イタリック
    p = tf.add_paragraph()
    p.text = "Italic text style"
    p.font.italic = True
    p.font.size = Pt(18)
    
    # 下線
    p = tf.add_paragraph()
    p.text = "Underlined text style"
    p.font.underline = True
    p.font.size = Pt(18)
    
    # 色付きテキスト
    p = tf.add_paragraph()
    p.text = "Colored text style"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 0, 0)  # 赤色
    
    # ファイルを保存
    prs.save('test_presentation.pptx')
    print("✅ テスト用PPTXファイルを作成しました: test_presentation.pptx")

if __name__ == "__main__":
    create_test_presentation()