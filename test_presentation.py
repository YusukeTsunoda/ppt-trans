#!/usr/bin/env python3
"""テスト用のPowerPointファイル作成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt

# プレゼンテーションを作成
prs = Presentation()

# スライド1: タイトルスライド
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Test Presentation"
subtitle.text = "This is a test PowerPoint file for upload testing"

# スライド2: コンテンツスライド
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Sample Content"
content.text = "• First bullet point\n• Second bullet point\n• Third bullet point"

# スライド3: 空白スライドにテキストボックス追加
slide_layout = prs.slide_layouts[5]  # 空白レイアウト
slide = prs.slides.add_slide(slide_layout)
left = Inches(1)
top = Inches(1)
width = Inches(8)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.text = "This is a custom text box"
text_frame.paragraphs[0].font.size = Pt(24)

# ファイルを保存
prs.save('test_presentation.pptx')
print("テスト用PowerPointファイル 'test_presentation.pptx' を作成しました。")