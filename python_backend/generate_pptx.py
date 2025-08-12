#!/usr/bin/env python3
"""
翻訳済みPPTXファイル生成スクリプト（改良版）
元のPPTXファイルと翻訳データを使って新しいPPTXファイルを生成する
フォント・レイアウトの保持機能とエラーハンドリングを強化
日本語フォント対応を改善
"""

import json
import sys
import io
import os
import traceback
import urllib.request
import tempfile
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import logging
import hashlib
from copy import deepcopy

# ログ設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# 日本語フォントのリスト（優先順位順）
JAPANESE_FONTS = [
    "游ゴシック",
    "Yu Gothic",
    "メイリオ",
    "Meiryo",
    "ＭＳ ゴシック",
    "MS Gothic",
    "ヒラギノ角ゴ Pro",
    "Hiragino Kaku Gothic Pro",
    "Noto Sans CJK JP",
    "源ノ角ゴシック",
]

# 英語フォントのリスト（フォールバック用）
ENGLISH_FONTS = [
    "Calibri",
    "Arial",
    "Helvetica",
    "Segoe UI",
]

@dataclass
class TextStyle:
    """テキストスタイル情報を保持するクラス"""
    font_name: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color_rgb: Optional[Tuple[int, int, int]] = None
    alignment: Optional[int] = None
    
    def apply_to_run(self, run, is_japanese: bool = False):
        """ランにスタイルを適用"""
        # フォントサイズを最優先で設定
        if self.font_size:
            run.font.size = self.font_size
            
        # フォント名の設定（日本語の場合は日本語フォントを優先）
        if is_japanese:
            # 元のフォントが日本語フォントならそれを維持
            if self.font_name and any(jp_font in self.font_name for jp_font in JAPANESE_FONTS):
                run.font.name = self.font_name
            else:
                # デフォルトの日本語フォントを設定
                run.font.name = JAPANESE_FONTS[0]
        elif self.font_name:
            run.font.name = self.font_name
            
        if self.bold is not None:
            run.font.bold = self.bold
        if self.italic is not None:
            run.font.italic = self.italic
        if self.underline is not None:
            run.font.underline = self.underline
        if self.color_rgb:
            try:
                run.font.color.rgb = RGBColor(*self.color_rgb)
            except:
                pass

@dataclass
class ParagraphStyle:
    """段落スタイル情報を保持するクラス"""
    alignment: Optional[int] = None
    level: Optional[int] = None
    line_spacing: Optional[float] = None
    space_before: Optional[int] = None
    space_after: Optional[int] = None
    
    def apply_to_paragraph(self, paragraph):
        """段落にスタイルを適用"""
        if self.alignment is not None:
            try:
                paragraph.alignment = self.alignment
            except:
                pass
        if self.level is not None:
            try:
                paragraph.level = self.level
            except:
                pass
        if self.line_spacing is not None:
            try:
                paragraph.line_spacing = self.line_spacing
            except:
                pass
        if self.space_before is not None:
            try:
                paragraph.space_before = self.space_before
            except:
                pass
        if self.space_after is not None:
            try:
                paragraph.space_after = self.space_after
            except:
                pass

class PPTXTranslator:
    """PPTXファイルの翻訳処理を行うクラス"""
    
    def __init__(self, original_file_path: str):
        """
        コンストラクタ
        
        Args:
            original_file_path: 元のPPTXファイルのパス（URLも可）
        """
        self.original_file_path = original_file_path
        self.presentation = None
        self.text_replacements = {}
        self.style_cache = {}
        self.error_log = []
        self.temp_file = None
        self.processed_shapes = set()  # 処理済みシェイプを追跡
        
    def download_if_url(self, file_path: str) -> str:
        """URLの場合はファイルをダウンロード"""
        if file_path.startswith(('http://', 'https://')):
            try:
                logger.info(f"Downloading file from URL: {file_path}")
                with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                    with urllib.request.urlopen(file_path) as response:
                        tmp_file.write(response.read())
                    self.temp_file = tmp_file.name
                    logger.info(f"Downloaded to temporary file: {self.temp_file}")
                    return self.temp_file
            except Exception as e:
                error_msg = f"Failed to download file from URL: {str(e)}"
                logger.error(error_msg)
                self.error_log.append(error_msg)
                raise
        return file_path
        
    def cleanup_temp_file(self):
        """一時ファイルをクリーンアップ"""
        if self.temp_file and os.path.exists(self.temp_file):
            try:
                os.unlink(self.temp_file)
                logger.info(f"Cleaned up temporary file: {self.temp_file}")
            except:
                pass
                
    def load_presentation(self) -> bool:
        """プレゼンテーションファイルを読み込む"""
        try:
            # URLの場合はダウンロード
            file_path = self.download_if_url(self.original_file_path)
            
            logger.info(f"Loading original PPTX: {file_path}")
            self.presentation = Presentation(file_path)
            logger.info(f"Successfully loaded presentation with {len(self.presentation.slides)} slides")
            return True
        except Exception as e:
            error_msg = f"Failed to load presentation: {str(e)}"
            logger.error(error_msg)
            self.error_log.append(error_msg)
            return False
    
    def is_japanese_text(self, text: str) -> bool:
        """テキストが日本語を含むか判定"""
        if not text:
            return False
        # 日本語文字の範囲をチェック
        for char in text:
            code = ord(char)
            # ひらがな、カタカナ、漢字の範囲
            if (0x3040 <= code <= 0x309F) or \
               (0x30A0 <= code <= 0x30FF) or \
               (0x4E00 <= code <= 0x9FFF) or \
               (0x3400 <= code <= 0x4DBF):
                return True
        return False
    
    def extract_text_style(self, run) -> TextStyle:
        """ランからテキストスタイルを抽出"""
        style = TextStyle()
        
        try:
            # フォント名を取得
            if run.font.name:
                style.font_name = run.font.name
            
            # フォントサイズを確実に取得（継承されている場合も考慮）
            if run.font.size:
                style.font_size = run.font.size
            elif hasattr(run, '_element'):
                # XMLから直接サイズを取得する試み
                try:
                    rPr = run._element.get_or_add_rPr()
                    sz = rPr.find('.//a:sz', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    if sz is not None and 'val' in sz.attrib:
                        # PowerPointのフォントサイズは100倍で保存されている
                        style.font_size = Pt(int(sz.attrib['val']) / 100)
                except:
                    pass
            
            if run.font.bold is not None:
                style.bold = run.font.bold
            if run.font.italic is not None:
                style.italic = run.font.italic
            if run.font.underline is not None:
                style.underline = run.font.underline
            
            # カラー情報の取得（エラーになる可能性があるため注意深く処理）
            try:
                if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                    style.color_rgb = (
                        run.font.color.rgb[0],
                        run.font.color.rgb[1],
                        run.font.color.rgb[2]
                    )
            except:
                pass
                
        except Exception as e:
            logger.debug(f"Could not extract some style properties: {e}")
            
        return style
    
    def extract_paragraph_style(self, paragraph) -> ParagraphStyle:
        """段落からスタイルを抽出"""
        style = ParagraphStyle()
        
        try:
            if hasattr(paragraph, 'alignment') and paragraph.alignment is not None:
                style.alignment = paragraph.alignment
            if hasattr(paragraph, 'level') and paragraph.level is not None:
                style.level = paragraph.level
            if hasattr(paragraph, 'line_spacing') and paragraph.line_spacing is not None:
                style.line_spacing = paragraph.line_spacing
            if hasattr(paragraph, 'space_before') and paragraph.space_before is not None:
                style.space_before = paragraph.space_before
            if hasattr(paragraph, 'space_after') and paragraph.space_after is not None:
                style.space_after = paragraph.space_after
        except Exception as e:
            logger.debug(f"Could not extract some paragraph properties: {e}")
            
        return style
    
    def prepare_text_replacements(self, edited_slides_data: List[Dict]) -> int:
        """
        翻訳マップを準備する
        
        Args:
            edited_slides_data: 編集済みスライドデータ
            
        Returns:
            準備された置換の数
        """
        self.text_replacements.clear()
        self.processed_shapes.clear()  # 処理済みシェイプをリセット
        
        for slide_data in edited_slides_data:
            for text_data in slide_data.get('texts', []):
                original = text_data.get('original', '').strip()
                translated = text_data.get('translated', '').strip()
                
                if original and translated and original != translated:
                    self.text_replacements[original] = translated
                    logger.debug(f"Added replacement: '{original[:50]}...' -> '{translated[:50]}...'")
        
        logger.info(f"Prepared {len(self.text_replacements)} text replacements")
        if logger.isEnabledFor(logging.DEBUG):
            for orig, trans in list(self.text_replacements.items())[:3]:
                logger.debug(f"Sample replacement: '{orig}' -> '{trans}'")
        
        return len(self.text_replacements)
    
    def replace_text_in_shape(self, shape, slide_idx: int, shape_idx: int) -> int:
        """
        シェイプ内のテキストを置換する（スタイル保持）
        
        Args:
            shape: 処理対象のシェイプ
            slide_idx: スライドインデックス
            shape_idx: シェイプインデックス
            
        Returns:
            置換されたテキストの数
        """
        if not shape.has_text_frame:
            return 0
        
        replaced_count = 0
        
        try:
            text_frame = shape.text_frame
            
            # テキストフレームの自動サイズ調整を無効化（サイズ保持のため）
            try:
                text_frame.auto_size = MSO_AUTO_SIZE.NONE
            except:
                pass
                
            for para_idx, paragraph in enumerate(text_frame.paragraphs):
                original_text = paragraph.text.strip()
                
                if original_text in self.text_replacements:
                    new_text = self.text_replacements[original_text]
                    is_japanese = self.is_japanese_text(new_text)
                    
                    logger.debug(f"Processing text shape paragraph {para_idx}: '{original_text}' -> '{new_text}'")
                    
                    # スタイルをキャッシュ
                    cache_key = f"{slide_idx}_{shape_idx}_{para_idx}"
                    
                    # 段落スタイルを保存
                    para_style = self.extract_paragraph_style(paragraph)
                    
                    # ランのスタイルを保存
                    run_styles = []
                    for run in paragraph.runs:
                        run_styles.append((run.text, self.extract_text_style(run)))
                    
                    self.style_cache[cache_key] = {
                        'paragraph': para_style,
                        'runs': run_styles
                    }
                    
                    # テキストを置換（重複を防ぐため一度だけ処理）
                    if paragraph.runs:
                        # 最初のランのスタイルを保持
                        first_run_style = self.extract_text_style(paragraph.runs[0])
                        
                        # すべてのランをクリアして新しいテキストを設定
                        paragraph.clear()
                        new_run = paragraph.add_run()
                        new_run.text = new_text
                        
                        # スタイルを適用
                        first_run_style.apply_to_run(new_run, is_japanese)
                        
                        # フォントサイズが設定されていない場合はデフォルトを設定
                        if not new_run.font.size and not first_run_style.font_size:
                            new_run.font.size = Pt(11)  # デフォルトサイズ
                            
                    else:
                        # ランがない場合は新しく作成
                        paragraph.text = new_text
                        if paragraph.runs and is_japanese:
                            # 日本語フォントを設定
                            paragraph.runs[0].font.name = JAPANESE_FONTS[0]
                    
                    # 段落スタイルを復元
                    para_style.apply_to_paragraph(paragraph)
                    
                    replaced_count += 1
                    logger.debug(f"Successfully replaced text in slide {slide_idx + 1}, shape {shape_idx + 1}, paragraph {para_idx + 1}")
                    
                    # 同じテキストの重複処理を防ぐため、処理済みをマークする
                    # (同一スライド内で同じoriginal_textの重複処理を避ける)
                    break
                    
        except Exception as e:
            error_msg = f"Error replacing text in shape (slide {slide_idx + 1}, shape {shape_idx + 1}): {str(e)}"
            logger.warning(error_msg)
            self.error_log.append(error_msg)
        
        return replaced_count
    
    def replace_text_in_table(self, shape, slide_idx: int) -> int:
        """テーブル内のテキストを置換"""
        if not shape.has_table:
            return 0
        
        replaced_count = 0
        
        try:
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    
                    if cell_text in self.text_replacements:
                        new_text = self.text_replacements[cell_text]
                        is_japanese = self.is_japanese_text(new_text)
                        
                        logger.debug(f"Processing table cell [{row_idx},{col_idx}]: '{cell_text}' -> '{new_text}'")
                        
                        # セルのテキストフレームが存在する場合の処理
                        if cell.text_frame and cell.text_frame.paragraphs:
                            # 最初の段落のみを使用して重複を防ぐ
                            first_para = cell.text_frame.paragraphs[0]
                            
                            # 元のスタイルを保存
                            original_style = None
                            if first_para.runs:
                                original_style = self.extract_text_style(first_para.runs[0])
                            
                            # 全ての段落をクリアしてから最初の段落にのみテキストを設定
                            for para in cell.text_frame.paragraphs:
                                para.clear()
                            
                            # 空の段落がある場合は削除（最初の段落は残す）
                            while len(cell.text_frame.paragraphs) > 1:
                                try:
                                    # 2番目以降の段落を削除
                                    cell.text_frame._element.remove(cell.text_frame.paragraphs[1]._element)
                                except:
                                    break
                            
                            # 最初の段落に新しいテキストを設定
                            new_run = first_para.add_run()
                            new_run.text = new_text
                            
                            # スタイルを適用
                            if original_style:
                                original_style.apply_to_run(new_run, is_japanese)
                            elif is_japanese:
                                new_run.font.name = JAPANESE_FONTS[0]
                                
                            # フォントサイズが未設定の場合はデフォルトを適用
                            if not new_run.font.size:
                                new_run.font.size = Pt(11)
                                
                        else:
                            # テキストフレームがない場合は直接設定
                            cell.text = new_text
                        
                        replaced_count += 1
                        logger.debug(f"Replaced text in table cell [{row_idx},{col_idx}] on slide {slide_idx + 1}")
                        
        except Exception as e:
            error_msg = f"Error replacing text in table (slide {slide_idx + 1}): {str(e)}"
            logger.warning(error_msg)
            self.error_log.append(error_msg)
        
        return replaced_count
    
    def process_slide(self, slide, slide_idx: int) -> int:
        """
        スライドを処理する
        
        Args:
            slide: 処理対象のスライド
            slide_idx: スライドインデックス
            
        Returns:
            置換されたテキストの数
        """
        replaced_count = 0
        
        logger.info(f"Processing slide {slide_idx + 1}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            # テキストフレームの処理
            replaced_count += self.replace_text_in_shape(shape, slide_idx, shape_idx)
            
            # テーブルの処理
            replaced_count += self.replace_text_in_table(shape, slide_idx)
            
            # グループシェイプの処理
            if shape.shape_type == 6:  # GROUP
                try:
                    for sub_shape in shape.shapes:
                        replaced_count += self.replace_text_in_shape(sub_shape, slide_idx, shape_idx)
                except:
                    pass
        
        return replaced_count
    
    def translate(self, edited_slides_data: List[Dict]) -> Tuple[bool, int]:
        """
        翻訳処理を実行する
        
        Args:
            edited_slides_data: 編集済みスライドデータ
            
        Returns:
            (成功フラグ, 置換されたテキストの総数)
        """
        try:
            # 翻訳マップを準備
            if self.prepare_text_replacements(edited_slides_data) == 0:
                logger.warning("No text replacements found")
                return True, 0
            
            # 各スライドを処理
            total_replaced = 0
            for slide_idx, slide in enumerate(self.presentation.slides):
                replaced = self.process_slide(slide, slide_idx)
                total_replaced += replaced
            
            logger.info(f"Total replacements: {total_replaced}")
            
            # スライドノートの処理（オプション）
            for slide_idx, slide in enumerate(self.presentation.slides):
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text in self.text_replacements:
                        slide.notes_slide.notes_text_frame.text = self.text_replacements[notes_text]
                        total_replaced += 1
                        logger.debug(f"Replaced notes text on slide {slide_idx + 1}")
            
            return True, total_replaced
            
        except Exception as e:
            error_msg = f"Translation failed: {str(e)}"
            logger.error(error_msg)
            self.error_log.append(error_msg)
            return False, 0
    
    def save(self, output_path: str) -> bool:
        """
        プレゼンテーションを保存する
        
        Args:
            output_path: 出力ファイルのパス
            
        Returns:
            成功した場合True
        """
        try:
            # 出力ディレクトリが存在しない場合は作成
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            logger.info(f"Saving translated PPTX to: {output_path}")
            self.presentation.save(output_path)
            logger.info("Translation completed successfully!")
            
            # 一時ファイルをクリーンアップ
            self.cleanup_temp_file()
            
            return True
        except Exception as e:
            error_msg = f"Failed to save presentation: {str(e)}"
            logger.error(error_msg)
            self.error_log.append(error_msg)
            return False
    
    def get_error_summary(self) -> str:
        """エラーサマリーを取得"""
        if not self.error_log:
            return "No errors"
        return "; ".join(self.error_log[:5])  # 最初の5つのエラーのみ

def generate_translated_pptx(
    original_file_path: str,
    edited_slides_data: List[Dict],
    output_path: str
) -> Dict[str, Any]:
    """
    翻訳済みPPTXファイルを生成する（メイン関数）
    
    Args:
        original_file_path: 元のPPTXファイルのパス（URLも可）
        edited_slides_data: 編集済みスライドデータ
        output_path: 出力ファイルのパス
    
    Returns:
        結果を含む辞書
    """
    result = {
        "success": False,
        "output": None,
        "replacements": 0,
        "errors": []
    }
    
    # URLの場合も処理可能
    if not original_file_path.startswith(('http://', 'https://')) and not os.path.exists(original_file_path):
        error_msg = f"Original file not found: {original_file_path}"
        logger.error(error_msg)
        result["errors"].append(error_msg)
        return result
    
    # トランスレーターを初期化
    translator = PPTXTranslator(original_file_path)
    
    # プレゼンテーションを読み込む
    if not translator.load_presentation():
        result["errors"] = translator.error_log
        return result
    
    # 翻訳処理を実行
    success, replacements = translator.translate(edited_slides_data)
    result["replacements"] = replacements
    
    if not success:
        result["errors"] = translator.error_log
        return result
    
    # ファイルを保存
    if translator.save(output_path):
        result["success"] = True
        result["output"] = output_path
        
        # 出力ファイルのサイズを確認
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            result["file_size"] = file_size
            logger.info(f"Output file size: {file_size:,} bytes")
    else:
        result["errors"] = translator.error_log
    
    # 警告やエラーがあれば追加
    if translator.error_log:
        result["warnings"] = translator.error_log
    
    return result

def main():
    """
    メイン処理
    コマンドライン引数：
        --input: 元のPPTXファイルパス（URLも可）
        --translations: 翻訳データJSONファイルパス
        --output: 出力ファイルパス
    """
    import argparse
    
    parser = argparse.ArgumentParser(description='Generate translated PPTX file')
    parser.add_argument('--input', required=True, help='Input PPTX file path or URL')
    parser.add_argument('--translations', required=True, help='Translation data JSON file path')
    parser.add_argument('--output', required=True, help='Output PPTX file path')
    
    args = parser.parse_args()
    
    try:
        # JSONデータを読み込む
        with open(args.translations, 'r', encoding='utf-8') as f:
            translation_data = json.load(f)
        
        # slidesキーがある場合はその中身を使用
        if 'slides' in translation_data:
            edited_slides = translation_data['slides']
        else:
            edited_slides = translation_data
        
        # PPTXファイルを生成
        result = generate_translated_pptx(args.input, edited_slides, args.output)
        
        # 結果を出力
        print(json.dumps(result, ensure_ascii=False, indent=2))
        
        # エラーコードを設定
        sys.exit(0 if result["success"] else 1)
        
    except json.JSONDecodeError as e:
        print(json.dumps({
            "success": False,
            "error": f"Invalid JSON: {str(e)}"
        }, ensure_ascii=False))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }, ensure_ascii=False))
        sys.exit(1)

if __name__ == "__main__":
    main()