#!/usr/bin/env python3
"""
翻訳済みのPPTXファイルを生成するスクリプト
元のPPTXファイルと翻訳データを基に、新しいPPTXファイルを作成します
"""

import argparse
import json
import sys
import os
from pathlib import Path
import tempfile
import shutil
from typing import Dict, List, Any
import traceback

# python-pptxパッケージをインポート
try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import MSO_AUTO_SIZE
except ImportError as e:
    # デバッグ情報を出力
    print(f"Error: python-pptx package is not installed. Import error: {e}", file=sys.stderr)
    print(f"Python executable: {sys.executable}", file=sys.stderr)
    print(f"Python path: {sys.path}", file=sys.stderr)
    sys.exit(1)

def read_translation_data(translations_path: str) -> Dict[str, Any]:
    """翻訳データをJSONファイルから読み込む"""
    try:
        with open(translations_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception as e:
        print(f"Error reading translation data: {e}", file=sys.stderr)
        raise

def download_file_from_url(url: str, output_path: str) -> None:
    """URLからファイルをダウンロード（Supabase対応）"""
    import urllib.request
    import urllib.parse
    
    try:
        # URLが相対パスの場合はローカルファイルとして処理
        if url.startswith('/'):
            # ローカルファイルパス
            local_path = url
            if url.startswith('/tmp/'):
                # /tmp/から始まる場合はそのまま使用
                local_path = url
            else:
                # プロジェクトルートからの相対パス
                local_path = os.path.join(os.getcwd(), url.lstrip('/'))
            
            if os.path.exists(local_path):
                shutil.copy2(local_path, output_path)
            else:
                raise FileNotFoundError(f"Local file not found: {local_path}")
        else:
            # HTTPまたはHTTPSのURL
            with urllib.request.urlopen(url) as response:
                with open(output_path, 'wb') as out_file:
                    shutil.copyfileobj(response, out_file)
    except Exception as e:
        print(f"Error downloading file from {url}: {e}", file=sys.stderr)
        raise

def replace_text_in_shape(shape, text_map: Dict[str, str]) -> None:
    """シェイプ内のテキストを置換（解決策A: Runレベルでの直接置換 + 解決策C: 自動調整制御）"""
    if not hasattr(shape, 'text_frame'):
        return
    
    text_frame = shape.text_frame
    if text_frame is None:
        return
    
    # 解決策C: 自動サイズ調整を無効化（テキストボックスのサイズ維持）
    original_auto_size = None
    try:
        original_auto_size = text_frame.auto_size
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except:
        # auto_sizeプロパティがない場合は無視
        pass
    
    # 解決策A: 段落ごとにRunレベルで直接置換
    for paragraph in text_frame.paragraphs:
        # 段落内の全てのrunのテキストを結合
        original_text = "".join(run.text for run in paragraph.runs).strip()
        
        # 空の段落はスキップ
        if not original_text:
            continue
        
        # テキストマップに存在する場合は置換
        if original_text in text_map:
            translated_text = text_map[original_text]
            
            # 最初のrunに翻訳済みテキストを設定（書式を維持）
            if paragraph.runs:
                # 最初のrunに全ての翻訳テキストを設定
                # これにより、フォントサイズ、色、太字などの書式が維持される
                paragraph.runs[0].text = translated_text
                
                # 2つ目以降のrunは空にする（書式の重複を避ける）
                for i in range(1, len(paragraph.runs)):
                    paragraph.runs[i].text = ""
            else:
                # runが存在しない場合は新規作成（稀なケース）
                run = paragraph.add_run()
                run.text = translated_text
    
    # 必要に応じて自動サイズ調整の設定を元に戻す（オプション）
    # ただし、通常は無効化したままの方が安定する
    # if original_auto_size is not None:
    #     text_frame.auto_size = original_auto_size

def process_slide(slide, slide_translations: List[Dict[str, Any]]) -> None:
    """スライド内のテキストを翻訳データで置換"""
    # 翻訳マップを作成（元のテキスト -> 翻訳後のテキスト）
    text_map = {}
    for text_data in slide_translations:
        original = text_data.get('original', '').strip()
        translated = text_data.get('translated', '').strip()
        if original and translated:
            text_map[original] = translated
    
    # スライド内のすべてのシェイプを処理
    for shape in slide.shapes:
        # テキストフレームを持つシェイプの処理
        replace_text_in_shape(shape, text_map)
        
        # テーブルの処理（同じ方法でテキストを置換）
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    # セルも同じ方法で処理
                    if hasattr(cell, 'text_frame') and cell.text_frame:
                        text_frame = cell.text_frame
                        
                        # 自動サイズ調整を無効化
                        try:
                            text_frame.auto_size = MSO_AUTO_SIZE.NONE
                        except:
                            pass
                        
                        # 段落ごとに処理
                        for paragraph in text_frame.paragraphs:
                            original_text = "".join(run.text for run in paragraph.runs).strip()
                            
                            if original_text and original_text in text_map:
                                translated_text = text_map[original_text]
                                
                                if paragraph.runs:
                                    paragraph.runs[0].text = translated_text
                                    for i in range(1, len(paragraph.runs)):
                                        paragraph.runs[i].text = ""
                                else:
                                    run = paragraph.add_run()
                                    run.text = translated_text
        
        # グループシェイプの処理
        if hasattr(shape, 'shapes'):
            for sub_shape in shape.shapes:
                replace_text_in_shape(sub_shape, text_map)

def generate_translated_pptx(input_path: str, translations_path: str, output_path: str) -> None:
    """翻訳済みPPTXファイルを生成"""
    try:
        # 一時ディレクトリを作成
        with tempfile.TemporaryDirectory() as temp_dir:
            # 元のPPTXファイルをダウンロード/コピー
            temp_input = os.path.join(temp_dir, 'input.pptx')
            download_file_from_url(input_path, temp_input)
            
            # 翻訳データを読み込む
            translation_data = read_translation_data(translations_path)
            slides_data = translation_data.get('slides', [])
            
            # PPTXファイルを開く
            prs = Presentation(temp_input)
            
            # 各スライドを処理
            for slide_data in slides_data:
                page_number = slide_data.get('pageNumber', 0)
                texts = slide_data.get('texts', [])
                
                # スライド番号は1から始まるが、配列インデックスは0から
                slide_index = page_number - 1
                
                if 0 <= slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    process_slide(slide, texts)
                else:
                    print(f"Warning: Slide {page_number} not found in presentation", file=sys.stderr)
            
            # 翻訳済みPPTXファイルを保存
            prs.save(output_path)
            
            print(f"Successfully generated translated PPTX: {output_path}")
            
    except Exception as e:
        print(f"Error generating translated PPTX: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr)
        raise

def main():
    """メイン処理"""
    parser = argparse.ArgumentParser(description='Generate translated PPTX file')
    parser.add_argument('--input', required=True, help='Input PPTX file path or URL')
    parser.add_argument('--translations', required=True, help='Translation data JSON file path')
    parser.add_argument('--output', required=True, help='Output PPTX file path')
    
    args = parser.parse_args()
    
    # 出力ディレクトリを作成
    output_dir = os.path.dirname(args.output)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    
    # 翻訳済みPPTXを生成
    generate_translated_pptx(args.input, args.translations, args.output)

if __name__ == '__main__':
    main()