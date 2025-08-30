#!/usr/bin/env python3
"""
テスト用のPowerPointファイルを生成するスクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import sys
import os

def create_test_pptx(output_path="test-presentation.pptx"):
    """
    テスト用のPowerPointファイルを作成
    
    Args:
        output_path: 出力ファイルパス
    """
    # プレゼンテーションを作成
    prs = Presentation()
    
    # スライド1: タイトルスライド
    slide_layout = prs.slide_layouts[0]  # タイトルスライドレイアウト
    slide1 = prs.slides.add_slide(slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Test Presentation"
    subtitle.text = "This is a test PowerPoint file for E2E testing"
    
    # スライド2: コンテンツスライド
    slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツ
    slide2 = prs.slides.add_slide(slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "Slide 2: Content"
    content.text = "• First bullet point\n• Second bullet point\n• Third bullet point"
    
    # スライド3: テーブルを含むスライド
    slide_layout = prs.slide_layouts[5]  # 空白レイアウト
    slide3 = prs.slides.add_slide(slide_layout)
    
    # タイトルを追加
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8)
    height = Inches(1)
    textbox = slide3.shapes.add_textbox(left, top, width, height)
    textbox.text_frame.text = "Slide 3: Table Example"
    
    # テーブルを追加
    rows = 3
    cols = 3
    left = Inches(2)
    top = Inches(2)
    width = Inches(5)
    height = Inches(3)
    
    table = slide3.shapes.add_table(rows, cols, left, top, width, height).table
    
    # テーブルにデータを追加
    table.cell(0, 0).text = "Header 1"
    table.cell(0, 1).text = "Header 2"
    table.cell(0, 2).text = "Header 3"
    table.cell(1, 0).text = "Row 1, Col 1"
    table.cell(1, 1).text = "Row 1, Col 2"
    table.cell(1, 2).text = "Row 1, Col 3"
    table.cell(2, 0).text = "Row 2, Col 1"
    table.cell(2, 1).text = "Row 2, Col 2"
    table.cell(2, 2).text = "Row 2, Col 3"
    
    # スライド4: 複数のテキストボックス
    slide_layout = prs.slide_layouts[5]  # 空白レイアウト
    slide4 = prs.slides.add_slide(slide_layout)
    
    # 複数のテキストボックスを追加
    textbox1 = slide4.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox1.text_frame.text = "Text Box 1"
    
    textbox2 = slide4.shapes.add_textbox(Inches(5), Inches(1), Inches(3), Inches(1))
    textbox2.text_frame.text = "Text Box 2"
    
    textbox3 = slide4.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(2))
    textbox3.text_frame.text = "This is a longer text that demonstrates how the extraction works with multiple lines of text in a single text box."
    
    # ファイルを保存
    prs.save(output_path)
    print(f"Test PowerPoint file created: {output_path}")
    return output_path

if __name__ == "__main__":
    output_path = sys.argv[1] if len(sys.argv) > 1 else "test-presentation.pptx"
    create_test_pptx(output_path)