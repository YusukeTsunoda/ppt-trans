#!/usr/bin/env python3
"""
翻訳済みPPTXファイル生成スクリプト（改良版）
元のPPTXファイルと翻訳データを使って新しいPPTXファイルを生成する
フォント・レイアウトの保持機能とエラーハンドリングを強化
"""

import json
import sys
import io
import os
import traceback
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import logging
import hashlib
from copy import deepcopy

# ログ設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

@dataclass
class TextStyle:
    """テキストスタイル情報を保持するクラス"""
    font_name: Optional[str] = None
    font_size: Optional[int] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None
    underline: Optional[bool] = None
    color_rgb: Optional[Tuple[int, int, int]] = None
    alignment: Optional[int] = None
    
    def apply_to_run(self, run):
        """ランにスタイルを適用"""
        if self.font_name:
            run.font.name = self.font_name
        if self.font_size:
            run.font.size = self.font_size
        if self.bold is not None:
            run.font.bold = self.bold
        if self.italic is not None:
            run.font.italic = self.italic
        if self.underline is not None:
            run.font.underline = self.underline
        if self.color_rgb:
            run.font.color.rgb = RGBColor(*self.color_rgb)

@dataclass
class ParagraphStyle:
    """段落スタイル情報を保持するクラス"""
    alignment: Optional[int] = None
    level: Optional[int] = None
    line_spacing: Optional[float] = None
    space_before: Optional[int] = None
    space_after: Optional[int] = None
    
    def apply_to_paragraph(self, paragraph):
        """段落にスタイルを適用"""
        if self.alignment is not None:
            paragraph.alignment = self.alignment
        if self.level is not None:
            paragraph.level = self.level
        if self.line_spacing is not None:
            paragraph.line_spacing = self.line_spacing
        if self.space_before is not None:
            paragraph.space_before = self.space_before
        if self.space_after is not None:
            paragraph.space_after = self.space_after

class PPTXTranslator:
    """PPTXファイルの翻訳処理を行うクラス"""
    
    def __init__(self, original_file_path: str):
        """
        コンストラクタ
        
        Args:
            original_file_path: 元のPPTXファイルのパス
        """
        self.original_file_path = original_file_path
        self.presentation = None
        self.text_replacements = {}
        self.style_cache = {}
        self.error_log = []
        
    def load_presentation(self) -> bool:
        """プレゼンテーションファイルを読み込む"""
        try:
            logger.info(f"Loading original PPTX: {self.original_file_path}")
            self.presentation = Presentation(self.original_file_path)
            logger.info(f"Successfully loaded presentation with {len(self.presentation.slides)} slides")
            return True
        except Exception as e:
            error_msg = f"Failed to load presentation: {str(e)}"
            logger.error(error_msg)
            self.error_log.append(error_msg)
            return False
    
    def extract_text_style(self, run) -> TextStyle:
        """ランからテキストスタイルを抽出"""
        style = TextStyle()
        
        try:
            if run.font.name:
                style.font_name = run.font.name
            if run.font.size:
                style.font_size = run.font.size
            if run.font.bold is not None:
                style.bold = run.font.bold
            if run.font.italic is not None:
                style.italic = run.font.italic
            if run.font.underline is not None:
                style.underline = run.font.underline
            
            # カラー情報の取得（エラーになる可能性があるため注意深く処理）
            try:
                if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                    style.color_rgb = (
                        run.font.color.rgb.r,
                        run.font.color.rgb.g,
                        run.font.color.rgb.b
                    )
            except:
                pass
                
        except Exception as e:
            logger.debug(f"Could not extract some style properties: {e}")
            
        return style
    
    def extract_paragraph_style(self, paragraph) -> ParagraphStyle:
        """段落からスタイルを抽出"""
        style = ParagraphStyle()
        
        try:
            if hasattr(paragraph, 'alignment') and paragraph.alignment is not None:
                style.alignment = paragraph.alignment
            if hasattr(paragraph, 'level') and paragraph.level is not None:
                style.level = paragraph.level
            if hasattr(paragraph, 'line_spacing') and paragraph.line_spacing is not None:
                style.line_spacing = paragraph.line_spacing
            if hasattr(paragraph, 'space_before') and paragraph.space_before is not None:
                style.space_before = paragraph.space_before
            if hasattr(paragraph, 'space_after') and paragraph.space_after is not None:
                style.space_after = paragraph.space_after
        except Exception as e:
            logger.debug(f"Could not extract some paragraph properties: {e}")
            
        return style
    
    def prepare_text_replacements(self, edited_slides_data: List[Dict]) -> int:
        """
        翻訳マップを準備する
        
        Args:
            edited_slides_data: 編集済みスライドデータ
            
        Returns:
            準備された置換の数
        """
        self.text_replacements.clear()
        
        for slide_data in edited_slides_data:
            for text_data in slide_data.get('texts', []):
                original = text_data.get('original', '').strip()
                translated = text_data.get('translated', '').strip()
                
                if original and translated and original != translated:
                    self.text_replacements[original] = translated
        
        logger.info(f"Prepared {len(self.text_replacements)} text replacements")
        return len(self.text_replacements)
    
    def replace_text_in_shape(self, shape, slide_idx: int, shape_idx: int) -> int:
        """
        シェイプ内のテキストを置換する（スタイル保持）
        
        Args:
            shape: 処理対象のシェイプ
            slide_idx: スライドインデックス
            shape_idx: シェイプインデックス
            
        Returns:
            置換されたテキストの数
        """
        if not shape.has_text_frame:
            return 0
        
        replaced_count = 0
        
        try:
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                original_text = paragraph.text.strip()
                
                if original_text in self.text_replacements:
                    new_text = self.text_replacements[original_text]
                    
                    # スタイルをキャッシュ
                    cache_key = f"{slide_idx}_{shape_idx}_{para_idx}"
                    
                    # 段落スタイルを保存
                    para_style = self.extract_paragraph_style(paragraph)
                    
                    # ランのスタイルを保存
                    run_styles = []
                    for run in paragraph.runs:
                        run_styles.append((run.text, self.extract_text_style(run)))
                    
                    self.style_cache[cache_key] = {
                        'paragraph': para_style,
                        'runs': run_styles
                    }
                    
                    # テキストを置換
                    if paragraph.runs:
                        # 日本語の場合、フォントを適切に設定
                        is_japanese = any(ord(char) > 0x3000 for char in new_text)
                        
                        # 最初のランのスタイルを保持
                        first_run_style = self.extract_text_style(paragraph.runs[0])
                        
                        # すべてのランをクリア
                        for i, run in enumerate(paragraph.runs):
                            if i == 0:
                                run.text = new_text
                                # スタイルを復元
                                first_run_style.apply_to_run(run)
                                
                                # 日本語フォントの設定
                                if is_japanese and not first_run_style.font_name:
                                    run.font.name = "游ゴシック"  # デフォルト日本語フォント
                            else:
                                run.text = ""
                    else:
                        # ランがない場合は新しく作成
                        paragraph.text = new_text
                        if paragraph.runs:
                            # 日本語フォントの設定
                            is_japanese = any(ord(char) > 0x3000 for char in new_text)
                            if is_japanese:
                                paragraph.runs[0].font.name = "游ゴシック"
                    
                    # 段落スタイルを復元
                    para_style.apply_to_paragraph(paragraph)
                    
                    replaced_count += 1
                    logger.debug(f"Replaced text in slide {slide_idx + 1}, shape {shape_idx + 1}")
                    
        except Exception as e:
            error_msg = f"Error replacing text in shape (slide {slide_idx + 1}, shape {shape_idx + 1}): {str(e)}"
            logger.warning(error_msg)
            self.error_log.append(error_msg)
        
        return replaced_count
    
    def replace_text_in_table(self, shape, slide_idx: int) -> int:
        """テーブル内のテキストを置換"""
        if not shape.has_table:
            return 0
        
        replaced_count = 0
        
        try:
            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text.strip()
                    
                    if cell_text in self.text_replacements:
                        new_text = self.text_replacements[cell_text]
                        
                        # セルのテキストフレームのスタイルを保持
                        if cell.text_frame:
                            for para in cell.text_frame.paragraphs:
                                if para.runs:
                                    # 最初のランのスタイルを保存
                                    style = self.extract_text_style(para.runs[0])
                                    para.text = new_text
                                    if para.runs:
                                        style.apply_to_run(para.runs[0])
                                else:
                                    para.text = new_text
                        else:
                            cell.text = new_text
                        
                        replaced_count += 1
                        logger.debug(f"Replaced text in table cell [{row_idx},{col_idx}] on slide {slide_idx + 1}")
                        
        except Exception as e:
            error_msg = f"Error replacing text in table (slide {slide_idx + 1}): {str(e)}"
            logger.warning(error_msg)
            self.error_log.append(error_msg)
        
        return replaced_count
    
    def process_slide(self, slide, slide_idx: int) -> int:
        """
        スライドを処理する
        
        Args:
            slide: 処理対象のスライド
            slide_idx: スライドインデックス
            
        Returns:
            置換されたテキストの数
        """
        replaced_count = 0
        
        logger.info(f"Processing slide {slide_idx + 1}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            # テキストフレームの処理
            replaced_count += self.replace_text_in_shape(shape, slide_idx, shape_idx)
            
            # テーブルの処理
            replaced_count += self.replace_text_in_table(shape, slide_idx)
            
            # グループシェイプの処理
            if shape.shape_type == 6:  # GROUP
                try:
                    for group_shape in shape.shapes:
                        replaced_count += self.replace_text_in_shape(group_shape, slide_idx, shape_idx)
                except:
                    pass  # グループシェイプへのアクセスに失敗した場合はスキップ
        
        return replaced_count
    
    def translate(self, edited_slides_data: List[Dict]) -> Tuple[bool, int]:
        """
        翻訳処理を実行
        
        Args:
            edited_slides_data: 編集済みスライドデータ
            
        Returns:
            (成功フラグ, 置換されたテキストの総数)
        """
        if not self.presentation:
            logger.error("Presentation not loaded")
            return False, 0
        
        # 翻訳マップを準備
        if self.prepare_text_replacements(edited_slides_data) == 0:
            logger.warning("No text replacements to apply")
            return True, 0
        
        # 各スライドを処理
        total_replaced = 0
        for slide_idx, slide in enumerate(self.presentation.slides):
            try:
                replaced = self.process_slide(slide, slide_idx)
                total_replaced += replaced
            except Exception as e:
                error_msg = f"Error processing slide {slide_idx + 1}: {str(e)}"
                logger.error(error_msg)
                self.error_log.append(error_msg)
                # エラーが発生してもスライドの処理を継続
        
        logger.info(f"Total replacements made: {total_replaced}")
        
        # 一部でも置換が成功していれば成功とする
        return total_replaced > 0 or len(self.text_replacements) == 0, total_replaced
    
    def save(self, output_path: str) -> bool:
        """
        プレゼンテーションを保存
        
        Args:
            output_path: 出力ファイルのパス
            
        Returns:
            成功した場合True
        """
        try:
            # 出力ディレクトリが存在しない場合は作成
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            logger.info(f"Saving translated PPTX to: {output_path}")
            self.presentation.save(output_path)
            logger.info("Translation completed successfully!")
            return True
        except Exception as e:
            error_msg = f"Failed to save presentation: {str(e)}"
            logger.error(error_msg)
            self.error_log.append(error_msg)
            return False
    
    def get_error_summary(self) -> str:
        """エラーサマリーを取得"""
        if not self.error_log:
            return "No errors"
        return "; ".join(self.error_log[:5])  # 最初の5つのエラーのみ

def generate_translated_pptx(
    original_file_path: str,
    edited_slides_data: List[Dict],
    output_path: str
) -> Dict[str, Any]:
    """
    翻訳済みPPTXファイルを生成する（メイン関数）
    
    Args:
        original_file_path: 元のPPTXファイルのパス
        edited_slides_data: 編集済みスライドデータ
        output_path: 出力ファイルのパス
    
    Returns:
        結果を含む辞書
    """
    result = {
        "success": False,
        "output": None,
        "replacements": 0,
        "errors": []
    }
    
    # 入力ファイルの存在確認
    if not os.path.exists(original_file_path):
        error_msg = f"Original file not found: {original_file_path}"
        logger.error(error_msg)
        result["errors"].append(error_msg)
        return result
    
    # トランスレーターを初期化
    translator = PPTXTranslator(original_file_path)
    
    # プレゼンテーションを読み込む
    if not translator.load_presentation():
        result["errors"] = translator.error_log
        return result
    
    # 翻訳処理を実行
    success, replacements = translator.translate(edited_slides_data)
    result["replacements"] = replacements
    
    if not success:
        result["errors"] = translator.error_log
        return result
    
    # ファイルを保存
    if translator.save(output_path):
        result["success"] = True
        result["output"] = output_path
        
        # 出力ファイルのサイズを確認
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            result["file_size"] = file_size
            logger.info(f"Output file size: {file_size:,} bytes")
    else:
        result["errors"] = translator.error_log
    
    # 警告やエラーがあれば追加
    if translator.error_log:
        result["warnings"] = translator.error_log
    
    return result

def main():
    """
    メイン処理
    コマンドライン引数：
        1. 元のPPTXファイルパス
        2. 編集済みスライドデータ（JSONファイルパスまたはJSON文字列）
        3. 出力ファイルパス
    """
    if len(sys.argv) != 4:
        print(json.dumps({
            "success": False,
            "error": "Usage: python generate_pptx_v2.py <original_pptx> <edited_slides_json> <output_pptx>"
        }))
        sys.exit(1)
    
    original_file = sys.argv[1]
    edited_slides_input = sys.argv[2]
    output_file = sys.argv[3]
    
    try:
        # JSONデータを読み込む（ファイルパスまたはJSON文字列）
        if os.path.exists(edited_slides_input):
            # JSONファイルとして読み込む
            with open(edited_slides_input, 'r', encoding='utf-8') as f:
                edited_slides = json.load(f)
        else:
            # JSON文字列として解析
            edited_slides = json.loads(edited_slides_input)
        
        # PPTXファイルを生成
        result = generate_translated_pptx(original_file, edited_slides, output_file)
        
        # 結果を出力
        print(json.dumps(result, ensure_ascii=False))
        
        # エラーコードを設定
        sys.exit(0 if result["success"] else 1)
        
    except json.JSONDecodeError as e:
        print(json.dumps({
            "success": False,
            "error": f"Invalid JSON: {str(e)}"
        }))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }))
        sys.exit(1)

if __name__ == "__main__":
    main()