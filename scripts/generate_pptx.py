#!/usr/bin/env python3
"""
翻訳済みPPTXファイル生成スクリプト
元のPPTXファイルと翻訳データを使って新しいPPTXファイルを生成する
"""

import json
import sys
import io
from pptx import Presentation
from pptx.util import Pt
import logging
import hashlib

# ログ設定
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def get_text_id(slide_idx, shape_idx, paragraph_idx, run_idx=0):
    """テキストの一意なIDを生成"""
    return hashlib.md5(f"{slide_idx}_{shape_idx}_{paragraph_idx}_{run_idx}".encode()).hexdigest()[:8]

def replace_text_in_shape(shape, text_replacements):
    """
    シェイプ内のテキストを置換する
    書式を可能な限り維持しながら置換を行う
    """
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        original_text = paragraph.text.strip()
        
        if original_text in text_replacements:
            new_text = text_replacements[original_text]
            
            # 元のランの書式を保存
            if paragraph.runs:
                # 最初のランの書式を保存
                first_run = paragraph.runs[0]
                font_name = first_run.font.name if first_run.font.name else None
                font_size = first_run.font.size if first_run.font.size else None
                font_bold = first_run.font.bold if first_run.font.bold else None
                font_italic = first_run.font.italic if first_run.font.italic else None
                font_color = first_run.font.color if hasattr(first_run.font, 'color') else None
                
                # すべてのランをクリア
                for run in paragraph.runs:
                    run.text = ""
                
                # 新しいテキストを最初のランに設定
                paragraph.runs[0].text = new_text
                
                # 書式を復元
                if font_name:
                    paragraph.runs[0].font.name = font_name
                if font_size:
                    paragraph.runs[0].font.size = font_size
                if font_bold is not None:
                    paragraph.runs[0].font.bold = font_bold
                if font_italic is not None:
                    paragraph.runs[0].font.italic = font_italic
            else:
                # ランがない場合は新しく作成
                paragraph.text = new_text

def generate_translated_pptx(original_file_path, edited_slides_data, output_path):
    """
    翻訳済みPPTXファイルを生成する
    
    Args:
        original_file_path: 元のPPTXファイルのパス
        edited_slides_data: 編集済みスライドデータ（JSON形式）
        output_path: 出力ファイルのパス
    
    Returns:
        bool: 成功した場合True
    """
    try:
        logger.info(f"Loading original PPTX: {original_file_path}")
        
        # 元のプレゼンテーションを読み込む
        pres = Presentation(original_file_path)
        
        # 翻訳マップを作成（元のテキスト -> 翻訳テキスト）
        text_replacements = {}
        for slide_data in edited_slides_data:
            for text_data in slide_data.get('texts', []):
                original = text_data.get('original', '').strip()
                translated = text_data.get('translated', '').strip()
                if original and translated:
                    text_replacements[original] = translated
        
        logger.info(f"Text replacements to apply: {len(text_replacements)}")
        
        # 各スライドのテキストを置換
        replaced_count = 0
        for slide_idx, slide in enumerate(pres.slides):
            logger.info(f"Processing slide {slide_idx + 1}")
            
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_text_frame:
                    for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        original_text = paragraph.text.strip()
                        
                        if original_text in text_replacements:
                            new_text = text_replacements[original_text]
                            logger.info(f"  Replacing: '{original_text[:50]}...' -> '{new_text[:50]}...'")
                            
                            # テキストを置換（書式を維持）
                            if paragraph.runs:
                                # 既存のランを使用
                                # すべてのランのテキストをクリアして、最初のランに新しいテキストを設定
                                for i, run in enumerate(paragraph.runs):
                                    if i == 0:
                                        run.text = new_text
                                    else:
                                        run.text = ""
                            else:
                                # ランがない場合は段落のテキストを直接設定
                                paragraph.text = new_text
                            
                            replaced_count += 1
                
                # テーブル内のテキストも処理
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text in text_replacements:
                                cell.text = text_replacements[cell_text]
                                replaced_count += 1
        
        logger.info(f"Total replacements made: {replaced_count}")
        
        # 新しいファイルとして保存
        logger.info(f"Saving translated PPTX to: {output_path}")
        pres.save(output_path)
        
        logger.info("Translation completed successfully!")
        return True
        
    except Exception as e:
        logger.error(f"Error generating translated PPTX: {str(e)}")
        return False

def main():
    """
    メイン処理
    コマンドライン引数：
        1. 元のPPTXファイルパス
        2. 編集済みスライドデータ（JSONファイルパス）
        3. 出力ファイルパス
    """
    if len(sys.argv) != 4:
        print("Usage: python generate_pptx.py <original_pptx> <edited_slides_json_file> <output_pptx>")
        sys.exit(1)
    
    original_file = sys.argv[1]
    edited_slides_json_file = sys.argv[2]
    output_file = sys.argv[3]
    
    try:
        # JSONファイルを読み込む
        with open(edited_slides_json_file, 'r', encoding='utf-8') as f:
            edited_slides = json.load(f)
        
        # PPTXファイルを生成
        success = generate_translated_pptx(original_file, edited_slides, output_file)
        
        if success:
            print(json.dumps({"success": True, "output": output_file}))
        else:
            print(json.dumps({"success": False, "error": "Failed to generate PPTX"}))
            sys.exit(1)
            
    except json.JSONDecodeError as e:
        print(json.dumps({"success": False, "error": f"Invalid JSON: {str(e)}"}))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()