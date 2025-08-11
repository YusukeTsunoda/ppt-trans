#!/usr/bin/env python3
"""
Create a simple test PowerPoint file
"""

from pptx import Presentation
from pptx.util import Inches

# Create presentation
prs = Presentation()

# Add title slide
slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Test Presentation"
subtitle.text = "This is a test presentation for translation"

# Add content slide
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "First Slide"
content.text = "This is the first slide content.\nIt has multiple lines.\nPlease translate this text."

# Add another slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Second Slide"
content.text = "This is the second slide.\nIt also needs translation.\nThank you for testing."

# Save the presentation
prs.save('test_presentation.pptx')
print("Test presentation created: test_presentation.pptx")