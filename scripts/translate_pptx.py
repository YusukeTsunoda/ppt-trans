#!/usr/bin/env python3
"""
PowerPoint翻訳スクリプト
"""

import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Any

# python-pptxのインポートチェック
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print(json.dumps({
        "error": "python-pptx is not installed. Please run: pip install python-pptx"
    }))
    sys.exit(1)

# Anthropic SDKのインポートチェック
try:
    from anthropic import Anthropic
except ImportError:
    print(json.dumps({
        "error": "anthropic is not installed. Please run: pip install anthropic"
    }))
    sys.exit(1)


def extract_text_from_slide(slide) -> List[str]:
    """スライドからテキストを抽出"""
    texts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text.strip())
        
        # テーブル内のテキストも抽出
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text:
                        texts.append(cell.text.strip())
    
    return [t for t in texts if t]  # 空でないテキストのみ返す


def translate_text(text: str, api_key: str, target_lang: str = "ja") -> str:
    """Anthropic APIを使用してテキストを翻訳"""
    try:
        client = Anthropic(api_key=api_key)
        
        prompt = f"""Translate the following text to {target_lang}. 
Keep the translation natural and appropriate for presentation slides.
Maintain any formatting like bullet points or numbering.
Only return the translated text, nothing else.

Text to translate:
{text}"""
        
        response = client.messages.create(
            model="claude-3-haiku-20240307",  # 高速で安価なモデルを使用
            max_tokens=1024,
            temperature=0.3,
            messages=[
                {"role": "user", "content": prompt}
            ]
        )
        
        return response.content[0].text
    except Exception as e:
        print(f"Translation error: {e}", file=sys.stderr)
        return text  # エラー時は元のテキストを返す


def update_slide_text(slide, translations: Dict[str, str]):
    """スライドのテキストを翻訳済みテキストで更新"""
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            original_text = shape.text.strip()
            if original_text in translations:
                # テキストフレームがある場合
                if hasattr(shape, "text_frame"):
                    shape.text_frame.clear()
                    shape.text = translations[original_text]
                else:
                    shape.text = translations[original_text]
        
        # テーブル内のテキストも更新
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text:
                        original_text = cell.text.strip()
                        if original_text in translations:
                            cell.text = translations[original_text]


def translate_presentation(input_path: str, output_path: str, api_key: str, target_lang: str = "ja") -> Dict[str, Any]:
    """プレゼンテーション全体を翻訳"""
    try:
        # プレゼンテーションを開く
        prs = Presentation(input_path)
        
        # すべてのユニークなテキストを収集
        all_texts = set()
        slide_count = 0
        
        for slide in prs.slides:
            slide_count += 1
            texts = extract_text_from_slide(slide)
            all_texts.update(texts)
        
        # テキストを翻訳
        translations = {}
        translated_count = 0
        
        for text in all_texts:
            if text:  # 空でないテキストのみ翻訳
                translated = translate_text(text, api_key, target_lang)
                translations[text] = translated
                translated_count += 1
                
                # 進捗を出力
                print(f"Translating: {translated_count}/{len(all_texts)}", file=sys.stderr)
        
        # スライドを更新
        for slide in prs.slides:
            update_slide_text(slide, translations)
        
        # 翻訳済みプレゼンテーションを保存
        prs.save(output_path)
        
        return {
            "success": True,
            "slide_count": slide_count,
            "text_count": len(all_texts),
            "translated_count": translated_count,
            "output_path": output_path
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": str(e)
        }


def main():
    """メイン処理"""
    if len(sys.argv) < 4:
        print(json.dumps({
            "error": "Usage: python translate_pptx.py <input_file> <output_file> <api_key> [target_lang]"
        }))
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    api_key = sys.argv[3]
    target_lang = sys.argv[4] if len(sys.argv) > 4 else "ja"
    
    # ファイルの存在確認
    if not os.path.exists(input_file):
        print(json.dumps({
            "error": f"Input file not found: {input_file}"
        }))
        sys.exit(1)
    
    # 出力ディレクトリの作成
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)
    
    # 翻訳実行
    result = translate_presentation(input_file, output_file, api_key, target_lang)
    
    # 結果を出力
    print(json.dumps(result, ensure_ascii=False, indent=2))
    
    # 成功/失敗でexit codeを設定
    sys.exit(0 if result.get("success") else 1)


if __name__ == "__main__":
    main()