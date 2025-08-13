#!/usr/bin/env python3
"""
PowerPoint to Image Converter using LibreOffice + pdf2image
This script converts PPTX files to preview images and extracts text information
"""

import sys
import json
import os
import subprocess
import tempfile
import uuid
from pathlib import Path
import io
import logging
from typing import List, Dict, Any

# Third-party imports
try:
    from pptx import Presentation
    from pdf2image import convert_from_path
    from PIL import Image
    from supabase import create_client, Client
    import requests
except ImportError as e:
    print(json.dumps({"error": f"Required Python packages not installed: {e}"}))
    sys.exit(1)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTXProcessor:
    def __init__(self, supabase_url: str, supabase_key: str):
        """Initialize the PPTX processor with Supabase client"""
        self.supabase: Client = create_client(supabase_url, supabase_key)
        
    def convert_pptx_to_pdf(self, pptx_path: str, temp_dir: str) -> str:
        """Convert PPTX to PDF using LibreOffice headless mode"""
        try:
            pdf_path = os.path.join(temp_dir, "converted.pdf")
            
            # LibreOffice command for headless conversion (use full path on macOS)
            libreoffice_cmd = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", temp_dir,
                pptx_path
            ]
            
            logger.info(f"Running LibreOffice conversion: {' '.join(libreoffice_cmd)}")
            
            # Execute LibreOffice conversion
            result = subprocess.run(
                libreoffice_cmd,
                capture_output=True,
                text=True,
                timeout=120  # 2-minute timeout
            )
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice conversion failed: {result.stderr}")
            
            # LibreOffice generates a PDF with the same base name as input
            input_name = Path(pptx_path).stem
            generated_pdf = os.path.join(temp_dir, f"{input_name}.pdf")
            
            if not os.path.exists(generated_pdf):
                raise Exception(f"Expected PDF file not found: {generated_pdf}")
                
            return generated_pdf
            
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timeout")
        except Exception as e:
            logger.error(f"PDF conversion error: {e}")
            raise
    
    def convert_pdf_to_images(self, pdf_path: str) -> List[Image.Image]:
        """Convert PDF to list of PIL Images using pdf2image"""
        try:
            logger.info(f"Converting PDF to images: {pdf_path}")
            
            # Convert PDF pages to images
            images = convert_from_path(
                pdf_path,
                dpi=150,  # Good balance between quality and file size
                fmt='JPEG',
                thread_count=2  # Optimize for performance
            )
            
            logger.info(f"Successfully converted PDF to {len(images)} images")
            return images
            
        except Exception as e:
            logger.error(f"PDF to image conversion error: {e}")
            raise Exception(f"Failed to convert PDF to images: {e}")
    
    def extract_text_from_slide(self, slide, page_number: int) -> List[Dict[str, Any]]:
        """Extract text and position information from a PowerPoint slide"""
        texts = []
        
        try:
            for shape in slide.shapes:
                # Check if shape has a table
                if shape.has_table:
                    table = shape.table
                    # Extract text from each cell in the table
                    for row_idx, row in enumerate(table.rows):
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            if not cell_text:
                                continue
                            
                            # Generate unique ID for this table cell text
                            text_id = f"table-{page_number}-r{row_idx}-c{col_idx}-{str(uuid.uuid4())[:8]}"
                            
                            # Extract position information for the table (use table's position)
                            try:
                                # Approximate cell position based on table position and cell indices
                                cell_width = float(shape.width.pt) / len(table.columns) if shape.width and len(table.columns) > 0 else 100.0
                                cell_height = float(shape.height.pt) / len(table.rows) if shape.height and len(table.rows) > 0 else 30.0
                                
                                position = {
                                    "x": float(shape.left.pt) + (col_idx * cell_width) if shape.left else col_idx * cell_width,
                                    "y": float(shape.top.pt) + (row_idx * cell_height) if shape.top else row_idx * cell_height,
                                    "width": cell_width,
                                    "height": cell_height
                                }
                            except:
                                # Fallback position if extraction fails
                                position = {"x": 0.0, "y": 0.0, "width": 100.0, "height": 30.0}
                            
                            text_data = {
                                "id": text_id,
                                "original": cell_text,
                                "translated": "",
                                "position": position,
                                "type": "table_cell"  # Mark as table cell for reference
                            }
                            
                            texts.append(text_data)
                            logger.info(f"Extracted table cell text: {cell_text[:50]}...")
                
                # Check if shape has a text frame (regular text)
                elif shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_content = paragraph.text.strip()
                        if not text_content:
                            continue
                        
                        # Generate unique ID for this text element
                        text_id = f"text-{page_number}-{str(uuid.uuid4())[:8]}"
                        
                        # Extract position information (convert from EMU to points)
                        try:
                            position = {
                                "x": float(shape.left.pt) if shape.left else 0.0,
                                "y": float(shape.top.pt) if shape.top else 0.0,
                                "width": float(shape.width.pt) if shape.width else 0.0,
                                "height": float(shape.height.pt) if shape.height else 0.0
                            }
                        except:
                            # Fallback position if extraction fails
                            position = {"x": 0.0, "y": 0.0, "width": 0.0, "height": 0.0}
                        
                        text_data = {
                            "id": text_id,
                            "original": text_content,
                            "translated": "",
                            "position": position,
                            "type": "text"  # Mark as regular text
                        }
                        
                        texts.append(text_data)
                    
        except Exception as e:
            logger.error(f"Text extraction error for slide {page_number}: {e}")
        
        # Sort texts by position: top to bottom, left to right
        # Define a threshold for grouping items on the same "row"
        row_threshold = 20  # points - items within 20 points vertically are considered same row
        
        if texts:
            # Sort by Y position first, then by X position
            # Group texts that are roughly on the same horizontal line
            texts.sort(key=lambda t: (
                round(t["position"]["y"] / row_threshold) * row_threshold,  # Group by row
                t["position"]["x"]  # Then sort by X within each row
            ))
            
            logger.info(f"Sorted {len(texts)} texts for slide {page_number} by position (top-to-bottom, left-to-right)")
            
        return texts
    
    def upload_image_to_supabase(self, image: Image.Image, filename: str) -> str:
        """Upload PIL Image to Supabase storage and return public URL"""
        try:
            # Convert PIL Image to bytes
            img_buffer = io.BytesIO()
            image.save(img_buffer, format='JPEG', quality=85, optimize=True)
            img_buffer.seek(0)
            
            # First, try to remove existing file if it exists
            try:
                self.supabase.storage.from_('pptx-files').remove([filename])
            except:
                pass  # Ignore errors if file doesn't exist
            
            # Upload to Supabase with upsert option
            upload_result = self.supabase.storage.from_('pptx-files').upload(
                path=filename,
                file=img_buffer.getvalue(),
                file_options={
                    "content-type": "image/jpeg",
                    "upsert": "true"  # Allow overwriting existing files
                }
            )
            
            # Check if upload was successful
            if hasattr(upload_result, 'error') and upload_result.error:
                raise Exception(f"Supabase upload error: {upload_result.error}")
            
            # Get public URL
            public_url_data = self.supabase.storage.from_('pptx-files').get_public_url(filename)
            
            # Extract the public URL
            if hasattr(public_url_data, 'public_url'):
                return public_url_data.public_url
            elif isinstance(public_url_data, dict) and 'publicUrl' in public_url_data:
                return public_url_data['publicUrl']
            elif isinstance(public_url_data, str):
                return public_url_data
            else:
                # If unable to get public URL, construct it manually
                return f"{self.supabase_url}/storage/v1/object/public/pptx-files/{filename}"
            
        except Exception as e:
            logger.error(f"Image upload error: {e}")
            raise Exception(f"Failed to upload image to Supabase: {e}")
    
    def process_pptx_file(self, pptx_path: str, temp_dir: str, temp_id: str) -> Dict[str, Any]:
        """Main processing function that coordinates the entire conversion pipeline"""
        try:
            logger.info(f"Starting PPTX processing: {pptx_path}")
            
            # Step 1: Load the PPTX file to extract text
            presentation = Presentation(pptx_path)
            logger.info(f"Loaded presentation with {len(presentation.slides)} slides")
            
            # Step 2: Convert PPTX to PDF using LibreOffice
            pdf_path = self.convert_pptx_to_pdf(pptx_path, temp_dir)
            
            # Step 3: Convert PDF to images
            images = self.convert_pdf_to_images(pdf_path)
            
            if len(images) != len(presentation.slides):
                logger.warning(f"Mismatch: {len(images)} images vs {len(presentation.slides)} slides")
            
            # Step 4: Process each slide
            slides_data = []
            
            for i, (slide, image) in enumerate(zip(presentation.slides, images)):
                page_number = i + 1
                logger.info(f"Processing slide {page_number}")
                
                # Extract text from slide
                texts = self.extract_text_from_slide(slide, page_number)
                
                # Upload image to Supabase
                image_filename = f"previews/{temp_id}_slide_{page_number}.jpg"
                image_url = self.upload_image_to_supabase(image, image_filename)
                
                slide_data = {
                    "pageNumber": page_number,
                    "imageUrl": image_url,
                    "texts": texts
                }
                
                slides_data.append(slide_data)
                logger.info(f"Completed slide {page_number} with {len(texts)} text elements")
            
            result = {
                "slides": slides_data,
                "totalSlides": len(slides_data)
            }
            
            logger.info(f"Successfully processed PPTX with {len(slides_data)} slides")
            return result
            
        except Exception as e:
            logger.error(f"PPTX processing error: {e}")
            raise

def main():
    """Main entry point for the script"""
    if len(sys.argv) != 4:
        print(json.dumps({"error": "Usage: python process_pptx.py <pptx_path> <temp_dir> <temp_id>"}))
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    temp_dir = sys.argv[2]
    temp_id = sys.argv[3]
    
    # Get environment variables
    supabase_url = os.getenv('SUPABASE_URL')
    supabase_key = os.getenv('SUPABASE_KEY')
    
    if not supabase_url or not supabase_key:
        print(json.dumps({"error": "SUPABASE_URL and SUPABASE_KEY environment variables are required"}))
        sys.exit(1)
    
    try:
        # Validate input file
        if not os.path.exists(pptx_path):
            raise Exception(f"Input file does not exist: {pptx_path}")
        
        # Initialize processor
        processor = PPTXProcessor(supabase_url, supabase_key)
        
        # Process the PPTX file
        result = processor.process_pptx_file(pptx_path, temp_dir, temp_id)
        
        # Output result as JSON
        print(json.dumps(result, ensure_ascii=False, indent=2))
        
    except Exception as e:
        error_result = {"error": str(e)}
        print(json.dumps(error_result))
        sys.exit(1)

if __name__ == "__main__":
    main()